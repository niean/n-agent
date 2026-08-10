"""Office artifact exporter (Infrastructure layer).

Implements the :class:`ArtifactExporter` Protocol using python-docx,
python-pptx, and openpyxl.  All third-party format libraries are confined
to this module (Infrastructure); the Application layer never imports them.

Security guarantees:
- No network access (images degrade to placeholder text).
- No local filesystem paths or usernames in Office document properties.
- XLSX formula-injection defence: cells whose first non-whitespace char
  is ``= + - @`` are forced to string type.
- No macros, external data connections, or formulas generated.
- Injectable size limits; violations raise
  :class:`ArtifactExportTooLargeError` without returning partial files.
"""
from __future__ import annotations

import io
import re
from dataclasses import dataclass
from typing import Mapping

from app.domain.artifact import (
    ArtifactExportError,
    ArtifactExportTooLargeError,
    ArtifactExportUnsupportedError,
    ArtifactKind,
    ArtifactRevision,
)
from app.domain.artifact_exporter import (
    ContentProfile,
    ExportedArtifact,
)

# MIME constants
_DOCX_MIME = (
    "application/vnd.openxmlformats-officedocument"
    ".wordprocessingml.document"
)
_PPTX_MIME = (
    "application/vnd.openxmlformats-officedocument"
    ".presentationml.presentation"
)
_XLSX_MIME = (
    "application/vnd.openxmlformats-officedocument"
    ".spreadsheetml.sheet"
)
_HTML_MIME = "text/html"

_MIB = 1024 * 1024

# Kinds treated as markdown-structured for capabilities.
_MARKDOWN_KINDS = frozenset({ArtifactKind.MARKDOWN, ArtifactKind.DOCUMENT})
_BINARY_KINDS = frozenset(
    {ArtifactKind.IMAGE, ArtifactKind.PDF, ArtifactKind.OTHER}
)


# ---------------------------------------------------------------------------
# Config
# ---------------------------------------------------------------------------


@dataclass(frozen=True)
class ArtifactExporterConfig:
    """Injectable export limits (all defaults align with spec)."""

    max_blocks: int = 10000
    max_table_cells: int = 100000
    max_slides: int = 200
    max_rows: int = 100000
    max_columns: int = 256
    max_cell_chars: int = 32767
    max_output_bytes: int = 50 * _MIB


# ---------------------------------------------------------------------------
# Markdown block parser
# ---------------------------------------------------------------------------


@dataclass
class _MdBlock:
    kind: str  # heading / paragraph / code / list_ol / list_ul / table / hr
    text: str = ""
    level: int = 0
    items: list[str] | None = None  # list items
    rows: list[list[str]] | None = None  # table rows


_HEADING_RE = re.compile(r"^(#{1,6})\s+(.*)$")
_OL_RE = re.compile(r"^(\d+)\.\s+(.*)$")
_UL_RE = re.compile(r"^[-*]\s+(.*)$")
_TABLE_SEP_RE = re.compile(r"^\|?[\s:|-]+\|?\s*$")
_BOM = "﻿"


def _parse_markdown(text: str) -> list[_MdBlock]:
    """Parse markdown into a flat list of structural blocks.

    Supports headings (H1-H6), fenced code blocks, ordered/unordered
    lists, pipe tables, horizontal rules, and paragraphs.  Inline HTML
    is stripped to visible text by :func:`_clean_inline`.
    """
    lines = text.split("\n")
    blocks: list[_MdBlock] = []
    i = 0
    n = len(lines)
    while i < n:
        line = lines[i]

        # skip blank lines
        if not line.strip():
            i += 1
            continue

        # fenced code block
        if line.lstrip().startswith("```"):
            code_lines: list[str] = []
            i += 1
            while i < n and not lines[i].lstrip().startswith("```"):
                code_lines.append(lines[i])
                i += 1
            i += 1  # skip closing fence
            blocks.append(_MdBlock(kind="code", text="\n".join(code_lines)))
            continue

        # heading
        m = _HEADING_RE.match(line)
        if m:
            level = len(m.group(1))
            content = _clean_inline(m.group(2).strip())
            blocks.append(
                _MdBlock(kind="heading", text=content, level=level)
            )
            i += 1
            continue

        # horizontal rule
        if re.match(r"^(-{3,}|\*{3,}|_{3,})$", line.strip()):
            blocks.append(_MdBlock(kind="hr"))
            i += 1
            continue

        # table: line starts with | and next line is a separator
        if "|" in line and i + 1 < n and _TABLE_SEP_RE.match(lines[i + 1].strip()):
            table_rows: list[list[str]] = []
            # header row
            table_rows.append(_split_table_row(line))
            i += 1  # skip header
            i += 1  # skip separator
            while i < n and "|" in lines[i] and lines[i].strip():
                table_rows.append(_split_table_row(lines[i]))
                i += 1
            blocks.append(_MdBlock(kind="table", rows=table_rows))
            continue

        # unordered list
        if _UL_RE.match(line):
            items: list[str] = []
            while i < n:
                m_ul = _UL_RE.match(lines[i])
                if m_ul:
                    items.append(_clean_inline(m_ul.group(1).strip()))
                    i += 1
                elif lines[i].strip() == "":
                    # allow blank line within list, check if next is list
                    if i + 1 < n and _UL_RE.match(lines[i + 1]):
                        i += 1
                    else:
                        break
                else:
                    break
            blocks.append(_MdBlock(kind="list_ul", items=items))
            continue

        # ordered list
        if _OL_RE.match(line):
            items_ol: list[str] = []
            while i < n:
                m_ol = _OL_RE.match(lines[i])
                if m_ol:
                    items_ol.append(_clean_inline(m_ol.group(2).strip()))
                    i += 1
                elif lines[i].strip() == "":
                    if i + 1 < n and _OL_RE.match(lines[i + 1]):
                        i += 1
                    else:
                        break
                else:
                    break
            blocks.append(_MdBlock(kind="list_ol", items=items_ol))
            continue

        # paragraph (collect consecutive non-blank lines)
        para_lines: list[str] = []
        while i < n and lines[i].strip() and not _is_block_start(lines[i], lines, i, n):
            para_lines.append(_clean_inline(lines[i].strip()))
            i += 1
        if para_lines:
            blocks.append(
                _MdBlock(kind="paragraph", text=" ".join(para_lines))
            )

    return blocks


def _is_block_start(
    line: str, lines: list[str], i: int, n: int
) -> bool:
    """Check if *line* begins a new block (heading, list, code, table, hr)."""
    if _HEADING_RE.match(line):
        return True
    if line.lstrip().startswith("```"):
        return True
    if re.match(r"^(-{3,}|\*{3,}|_{3,})$", line.strip()):
        return True
    if _UL_RE.match(line):
        return True
    if _OL_RE.match(line):
        return True
    if "|" in line and i + 1 < n and _TABLE_SEP_RE.match(lines[i + 1].strip()):
        return True
    return False


def _split_table_row(line: str) -> list[str]:
    """Split a markdown table row into cell values."""
    s = line.strip()
    if s.startswith("|"):
        s = s[1:]
    if s.endswith("|"):
        s = s[:-1]
    return [_clean_inline(c.strip()) for c in s.split("|")]


# ---------------------------------------------------------------------------
# Inline text cleaning
# ---------------------------------------------------------------------------

_IMG_RE = re.compile(r"!\[([^\]]*)\]\([^)]*\)")
_LINK_RE = re.compile(r"\[([^\]]*)\]\([^)]*\)")
_BOLD_RE = re.compile(r"\*\*(.+?)\*\*")
_ITALIC_RE = re.compile(r"\*(.+?)\*")
_CODE_RE = re.compile(r"`([^`]+)`")
_HTML_TAG_RE = re.compile(r"<[^>]+>")


def _clean_inline(text: str) -> str:
    """Strip markdown inline formatting and degrade HTML to text.

    Images become ``[Image: alt]`` placeholders; links keep their text;
    HTML tags are removed (content preserved).  No network access.
    """
    # images -> placeholder (preserve alt text)
    text = _IMG_RE.sub(lambda m: f"[Image: {m.group(1)}]" if m.group(1) else "[Image]", text)
    # links -> text only
    text = _LINK_RE.sub(r"\1", text)
    # bold
    text = _BOLD_RE.sub(r"\1", text)
    # italic
    text = _ITALIC_RE.sub(r"\1", text)
    # inline code
    text = _CODE_RE.sub(r"\1", text)
    # strip raw HTML tags (degrade to text content)
    text = _HTML_TAG_RE.sub("", text)
    return text.strip()


# ---------------------------------------------------------------------------
# Safe filename stem
# ---------------------------------------------------------------------------


_UNSAFE_CHARS_RE = re.compile(r"[^a-zA-Z0-9._-]")


def _safe_stem(name: str) -> str:
    """Return a safe filename stem from *name*.

    Strips path components, removes the extension, and replaces
    unsafe characters with underscores.
    """
    # take basename only
    name = name.replace("\\", "/").split("/")[-1]
    # remove extension
    if "." in name:
        name = name.rsplit(".", 1)[0]
    # replace unsafe chars
    name = _UNSAFE_CHARS_RE.sub("_", name)
    name = name.strip("._-")
    return name or "artifact"


# ---------------------------------------------------------------------------
# XLSX cell value helpers
# ---------------------------------------------------------------------------

_DECIMAL_RE = re.compile(r"^-?\d+(\.\d+)?$")
_FORMULA_PREFIXES = frozenset({"=", "+", "-", "@"})


def _csv_cell_value(raw: str, cfg: ArtifactExporterConfig, col_idx: int):
    """Return (value, data_type) for a CSV cell.

    Formula-injection defence: after BOM removal, if the first
    non-whitespace char is ``= + - @``, the cell is forced to string.
    Strict-decimal fields (without leading ``+``) that pass the injection
    check become numbers.
    """
    # strip BOM
    cell = raw.replace(_BOM, "")
    if len(cell) > cfg.max_cell_chars:
        raise ArtifactExportTooLargeError(
            f"cell exceeds max_cell_chars ({cfg.max_cell_chars})"
        )
    stripped = cell.lstrip()
    if stripped and stripped[0] in _FORMULA_PREFIXES:
        # force string to prevent formula injection
        return cell
    # strict decimal (no leading +) -> number
    check = cell.strip()
    if check and not check.startswith("+") and _DECIMAL_RE.match(check):
        try:
            if "." in check:
                return float(check)
            return int(check)
        except ValueError:
            pass
    return cell


def _json_cell_value(value, cfg: ArtifactExporterConfig):
    """Return the xlsx value for a JSON-sourced cell."""
    if value is None:
        return None
    if isinstance(value, bool):
        return value
    if isinstance(value, (int, float)):
        return value
    if isinstance(value, str):
        if len(value) > cfg.max_cell_chars:
            raise ArtifactExportTooLargeError(
                f"cell exceeds max_cell_chars ({cfg.max_cell_chars})"
            )
        # formula injection defence for string values
        stripped = value.lstrip()
        if stripped and stripped[0] in _FORMULA_PREFIXES:
            return value
        return value
    # nested objects/arrays shouldn't reach here (filtered by probe)
    return str(value)


# ---------------------------------------------------------------------------
# Exporter
# ---------------------------------------------------------------------------


class OfficeArtifactExporter:
    """Infrastructure exporter producing DOCX/PPTX/XLSX/HTML/original.

    Implements the :class:`ArtifactExporter` Protocol.  All format
    libraries are imported lazily inside the export methods so that
    importing this module does not require them at collection time.
    """

    def __init__(self, config: ArtifactExporterConfig | None = None) -> None:
        self._cfg = config or ArtifactExporterConfig()

    # -- capabilities ----------------------------------------------------

    async def capabilities(
        self,
        kind: ArtifactKind,
        mime: str,
        content_profile: ContentProfile,
    ) -> tuple[str, ...]:
        fmts: set[str] = {"original"}
        if kind in _MARKDOWN_KINDS and content_profile is ContentProfile.MARKDOWN:
            fmts.update({"docx", "html", "pptx"})
        elif kind is ArtifactKind.TEXT and content_profile is ContentProfile.TEXT:
            fmts.add("docx")
        elif (
            kind in (ArtifactKind.CSV, ArtifactKind.JSON, ArtifactKind.DATA)
            and content_profile is ContentProfile.TABULAR
        ):
            fmts.add("xlsx")
        return tuple(sorted(fmts))

    # -- export ----------------------------------------------------------

    async def export(
        self,
        revision: ArtifactRevision,
        content_bytes: bytes,
        format: str,
        options: Mapping[str, object] | None = None,
    ) -> ExportedArtifact:
        fmt = format.lower()
        opts = options or {}
        content_profile = opts.get("content_profile")
        if not isinstance(content_profile, ContentProfile):
            content_profile = ContentProfile.TEXT
        artifact_name = str(opts.get("artifact_name", "artifact"))
        stem = _safe_stem(artifact_name)

        if fmt == "original":
            return self._to_original(revision, content_bytes, artifact_name)
        if fmt == "html":
            return self._to_html(content_bytes, stem)
        if fmt == "docx":
            return self._to_docx(revision, content_bytes, content_profile, stem)
        if fmt == "pptx":
            return self._to_pptx(revision, content_bytes, content_profile, stem)
        if fmt == "xlsx":
            return self._to_xlsx(revision, content_bytes, content_profile, stem)
        raise ArtifactExportUnsupportedError(
            f"unsupported export format: {fmt}"
        )

    # -- original --------------------------------------------------------

    def _to_original(
        self,
        revision: ArtifactRevision,
        content_bytes: bytes,
        artifact_name: str,
    ) -> ExportedArtifact:
        if not content_bytes:
            raise ArtifactExportError("cannot export empty content as original")
        self._check_output_size(content_bytes)
        # keep original filename (sanitised but preserving extension)
        safe = _safe_filename(artifact_name)
        return ExportedArtifact(
            data=content_bytes,
            mime=revision.mime,
            filename=safe,
        )

    # -- html ------------------------------------------------------------

    def _to_html(self, content_bytes: bytes, stem: str) -> ExportedArtifact:
        from app.infrastructure.artifact.export_converter import convert_to_html

        text = content_bytes.decode("utf-8", errors="replace")
        html = convert_to_html(text)
        data = html.encode("utf-8")
        self._check_output_size(data)
        return ExportedArtifact(data=data, mime=_HTML_MIME, filename=f"{stem}.html")

    # -- docx ------------------------------------------------------------

    def _to_docx(
        self,
        revision: ArtifactRevision,
        content_bytes: bytes,
        profile: ContentProfile,
        stem: str,
    ) -> ExportedArtifact:
        from docx import Document
        from docx.shared import Pt

        kind = revision.kind
        # DOCX supports markdown/document (MARKDOWN) and TEXT kind (TEXT).
        if profile is ContentProfile.MARKDOWN and kind in _MARKDOWN_KINDS:
            text = content_bytes.decode("utf-8", errors="replace")
            blocks = _parse_markdown(text)
        elif kind is ArtifactKind.TEXT and profile is ContentProfile.TEXT:
            text = content_bytes.decode("utf-8", errors="replace")
            blocks = [
                _MdBlock(kind="paragraph", text=line)
                for line in text.split("\n")
                if line.strip()
            ]
        else:
            raise ArtifactExportUnsupportedError(
                f"docx export not supported for kind={kind.value}, "
                f"profile={profile.value}"
            )

        if len(blocks) > self._cfg.max_blocks:
            raise ArtifactExportTooLargeError(
                f"markdown blocks ({len(blocks)}) exceed max_blocks "
                f"({self._cfg.max_blocks})"
            )

        doc = Document()
        # scrub document properties
        cp = doc.core_properties
        cp.author = ""
        cp.last_modified_by = ""

        for block in blocks:
            if block.kind == "heading":
                doc.add_heading(block.text, level=min(block.level, 6))
            elif block.kind == "paragraph":
                doc.add_paragraph(block.text)
            elif block.kind == "code":
                p = doc.add_paragraph(block.text)
                if p.runs:
                    for run in p.runs:
                        run.font.name = "Courier New"
                        run.font.size = Pt(10)
            elif block.kind == "list_ul":
                for item in block.items or []:
                    doc.add_paragraph(item, style="List Bullet")
            elif block.kind == "list_ol":
                for item in block.items or []:
                    doc.add_paragraph(item, style="List Number")
            elif block.kind == "table":
                self._docx_add_table(doc, block.rows or [])
            elif block.kind == "hr":
                doc.add_paragraph("---")

        buf = io.BytesIO()
        doc.save(buf)
        data = buf.getvalue()
        self._check_output_size(data)
        return ExportedArtifact(data=data, mime=_DOCX_MIME, filename=f"{stem}.docx")

    def _docx_add_table(self, doc, rows: list[list[str]]) -> None:
        if not rows:
            return
        n_cols = max(len(r) for r in rows)
        total_cells = len(rows) * n_cols
        if total_cells > self._cfg.max_table_cells:
            raise ArtifactExportTooLargeError(
                f"table cells ({total_cells}) exceed max_table_cells "
                f"({self._cfg.max_table_cells})"
            )
        table = doc.add_table(rows=len(rows), cols=n_cols)
        for i, row in enumerate(rows):
            for j in range(n_cols):
                cell_text = row[j] if j < len(row) else ""
                table.rows[i].cells[j].text = cell_text

    # -- pptx ------------------------------------------------------------

    def _to_pptx(
        self,
        revision: ArtifactRevision,
        content_bytes: bytes,
        profile: ContentProfile,
        stem: str,
    ) -> ExportedArtifact:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        kind = revision.kind
        # PPTX supports markdown/document (MARKDOWN) only.
        if not (profile is ContentProfile.MARKDOWN and kind in _MARKDOWN_KINDS):
            raise ArtifactExportUnsupportedError(
                f"pptx export not supported for kind={kind.value}, "
                f"profile={profile.value}"
            )

        text = content_bytes.decode("utf-8", errors="replace")
        blocks = _parse_markdown(text)

        # count slides (each H2 -> one slide; H1 -> title slide)
        slide_count = 0
        for b in blocks:
            if b.kind == "heading":
                if b.level == 1:
                    slide_count += 1
                elif b.level == 2:
                    slide_count += 1
        if slide_count == 0:
            slide_count = 1  # at least one slide
        if slide_count > self._cfg.max_slides:
            raise ArtifactExportTooLargeError(
                f"slides ({slide_count}) exceed max_slides "
                f"({self._cfg.max_slides})"
            )

        if len(blocks) > self._cfg.max_blocks:
            raise ArtifactExportTooLargeError(
                f"markdown blocks ({len(blocks)}) exceed max_blocks "
                f"({self._cfg.max_blocks})"
            )

        prs = Presentation()
        # scrub properties
        pcp = prs.core_properties
        pcp.author = ""
        pcp.last_modified_by = ""

        blank_layout = prs.slide_layouts[6]  # blank
        title_layout = prs.slide_layouts[0]  # title slide

        current_slide = None
        body_texts: list[str] = []

        def _flush_body(slide, texts):
            if not texts:
                return
            left = top = Inches(1)
            width = Inches(8)
            height = Inches(5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            for idx, t in enumerate(texts):
                p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                p.text = t

        for block in blocks:
            if block.kind == "heading" and block.level == 1:
                # flush previous slide body
                if current_slide is not None:
                    _flush_body(current_slide, body_texts)
                    body_texts = []
                slide = prs.slides.add_slide(title_layout)
                if slide.shapes.title:
                    slide.shapes.title.text = block.text
                current_slide = slide
            elif block.kind == "heading" and block.level == 2:
                if current_slide is not None:
                    _flush_body(current_slide, body_texts)
                    body_texts = []
                slide = prs.slides.add_slide(blank_layout)
                # add title textbox
                left = top = Inches(0.5)
                width = Inches(9)
                height = Inches(1)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.text = block.text
                if tf.paragraphs[0].runs:
                    tf.paragraphs[0].runs[0].font.size = Pt(28)
                    tf.paragraphs[0].runs[0].font.bold = True
                current_slide = slide
            elif block.kind == "heading":
                # H3-H6 -> body text
                body_texts.append(block.text)
            elif block.kind == "paragraph":
                body_texts.append(block.text)
            elif block.kind == "code":
                body_texts.append(block.text)
            elif block.kind == "list_ul":
                for item in block.items or []:
                    body_texts.append(f"- {item}")
            elif block.kind == "list_ol":
                for idx, item in enumerate(block.items or [], 1):
                    body_texts.append(f"{idx}. {item}")
            elif block.kind == "table":
                if current_slide is not None:
                    self._pptx_add_table(current_slide, block.rows or [])
            elif block.kind == "hr":
                body_texts.append("---")

        # flush last slide
        if current_slide is not None:
            _flush_body(current_slide, body_texts)
        elif current_slide is None:
            # no H1/H2 at all -> single slide with all body text
            slide = prs.slides.add_slide(blank_layout)
            all_texts = []
            for block in blocks:
                if block.kind == "paragraph":
                    all_texts.append(block.text)
                elif block.kind == "list_ul":
                    all_texts.extend(f"- {it}" for it in block.items or [])
                elif block.kind == "list_ol":
                    all_texts.extend(
                        f"{i}. {it}" for i, it in enumerate(block.items or [], 1)
                    )
            _flush_body(slide, all_texts)

        buf = io.BytesIO()
        prs.save(buf)
        data = buf.getvalue()
        self._check_output_size(data)
        return ExportedArtifact(data=data, mime=_PPTX_MIME, filename=f"{stem}.pptx")

    def _pptx_add_table(self, slide, rows: list[list[str]]) -> None:
        from pptx.util import Inches

        if not rows:
            return
        n_cols = max(len(r) for r in rows)
        total_cells = len(rows) * n_cols
        if total_cells > self._cfg.max_table_cells:
            raise ArtifactExportTooLargeError(
                f"table cells ({total_cells}) exceed max_table_cells "
                f"({self._cfg.max_table_cells})"
            )
        left = top = Inches(2)
        width = Inches(6)
        height = Inches(2)
        table_shape = slide.shapes.add_table(
            len(rows), n_cols, left, top, width, height
        )
        table = table_shape.table
        for i, row in enumerate(rows):
            for j in range(n_cols):
                cell_text = row[j] if j < len(row) else ""
                table.cell(i, j).text = cell_text

    # -- xlsx ------------------------------------------------------------

    def _to_xlsx(
        self,
        revision: ArtifactRevision,
        content_bytes: bytes,
        profile: ContentProfile,
        stem: str,
    ) -> ExportedArtifact:
        import openpyxl

        if profile is not ContentProfile.TABULAR:
            raise ArtifactExportUnsupportedError(
                "xlsx export requires tabular content profile"
            )

        # parse content into rows
        rows: list[list] = self._parse_tabular(revision, content_bytes)
        if not rows:
            raise ArtifactExportError("no tabular data to export")

        n_rows = len(rows)
        n_cols = max(len(r) for r in rows)
        if n_rows > self._cfg.max_rows:
            raise ArtifactExportTooLargeError(
                f"rows ({n_rows}) exceed max_rows ({self._cfg.max_rows})"
            )
        if n_cols > self._cfg.max_columns:
            raise ArtifactExportTooLargeError(
                f"columns ({n_cols}) exceed max_columns ({self._cfg.max_columns})"
            )

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Artifact"
        # scrub properties
        wb.properties.creator = ""
        wb.properties.lastModifiedBy = ""
        wb.properties.title = ""
        wb.properties.subject = ""
        wb.properties.description = ""

        kind = revision.kind
        for i, row in enumerate(rows):
            for j in range(len(row)):
                raw = row[j]
                if kind is ArtifactKind.CSV:
                    value = _csv_cell_value(str(raw), self._cfg, j)
                else:
                    value = _json_cell_value(raw, self._cfg)
                cell = ws.cell(row=i + 1, column=j + 1)
                cell.value = value
                # Force string type for dangerous prefixes to prevent
                # openpyxl from auto-detecting them as formulas.
                if (
                    isinstance(value, str)
                    and value
                    and value.lstrip()
                    and value.lstrip()[0] in _FORMULA_PREFIXES
                ):
                    cell.data_type = "s"

        buf = io.BytesIO()
        wb.save(buf)
        data = buf.getvalue()
        self._check_output_size(data)
        return ExportedArtifact(data=data, mime=_XLSX_MIME, filename=f"{stem}.xlsx")

    def _parse_tabular(
        self, revision: ArtifactRevision, content_bytes: bytes
    ) -> list[list]:
        """Parse CSV or JSON content into a list of rows."""
        import csv as csv_mod
        import json as json_mod

        kind = revision.kind
        if kind is ArtifactKind.CSV:
            text = content_bytes.decode("utf-8", errors="replace")
            reader = csv_mod.reader(io.StringIO(text))
            return [list(row) for row in reader]
        if kind in (ArtifactKind.JSON, ArtifactKind.DATA):
            text = content_bytes.decode("utf-8", errors="replace")
            parsed = json_mod.loads(text)
            if isinstance(parsed, list) and parsed:
                if all(isinstance(el, list) for el in parsed):
                    return [list(el) for el in parsed]
                if all(isinstance(el, dict) for el in parsed):
                    keys = list(parsed[0].keys())
                    result = [list(keys)]
                    for obj in parsed:
                        result.append([obj.get(k) for k in keys])
                    return result
        # fallback: try CSV
        text = content_bytes.decode("utf-8", errors="replace")
        reader = csv_mod.reader(io.StringIO(text))
        return [list(row) for row in reader]

    # -- helpers ---------------------------------------------------------

    def _check_output_size(self, data: bytes) -> None:
        if not data:
            raise ArtifactExportError("export produced empty output")
        if len(data) > self._cfg.max_output_bytes:
            raise ArtifactExportTooLargeError(
                f"output size ({len(data)}) exceeds max_output_bytes "
                f"({self._cfg.max_output_bytes})"
            )


def _safe_filename(name: str) -> str:
    """Return a safe filename (preserving extension)."""
    name = name.replace("\\", "/").split("/")[-1]
    name = _UNSAFE_CHARS_RE.sub("_", name)
    name = name.strip("._-")
    return name or "artifact"
