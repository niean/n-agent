"""Tests for OfficeArtifactExporter (Infrastructure layer).

Covers DOCX/PPTX/XLSX/HTML/original export, capabilities matrix, formula
injection defence, metadata scrubbing, injectable limits, and markdown
structure round-tripping via the corresponding parse libraries.
"""
from __future__ import annotations

import hashlib
import io
import socket
import types
import zipfile
from contextlib import contextmanager
from datetime import datetime, timezone
from unittest.mock import patch

import pytest

from app.domain.artifact import (
    ArtifactExportError,
    ArtifactExportTooLargeError,
    ArtifactExportUnsupportedError,
    ArtifactKind,
    ArtifactRevision,
)
from app.domain.artifact_exporter import ContentProfile, ExportedArtifact
from app.infrastructure.artifact.exporters import (
    ArtifactExporterConfig,
    OfficeArtifactExporter,
)

_TEXT_KINDS = frozenset({
    ArtifactKind.DOCUMENT,
    ArtifactKind.MARKDOWN,
    ArtifactKind.CODE,
    ArtifactKind.HTML,
    ArtifactKind.DATA,
    ArtifactKind.CSV,
    ArtifactKind.JSON,
    ArtifactKind.TEXT,
})

_DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
_XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
_HTML_MIME = "text/html"


# ---------------------------------------------------------------------------
# Test helpers
# ---------------------------------------------------------------------------


def _rev(
    data: bytes, *, kind: ArtifactKind, mime: str
) -> tuple[ArtifactRevision, bytes]:
    """Build a legal ArtifactRevision + content bytes."""
    checksum = "sha256:" + hashlib.sha256(data).hexdigest()
    now = datetime.now(timezone.utc)
    if kind in _TEXT_KINDS:
        rev = ArtifactRevision(
            id="rev1",
            artifact_id="art1",
            revision_number=1,
            parent_revision_id=None,
            rollback_from_revision_id=None,
            content_ref=None,
            inline_content=data.decode("utf-8"),
            size=len(data),
            checksum=checksum,
            mime=mime,
            kind=kind,
            created_at=now,
        )
    else:
        rev = ArtifactRevision(
            id="rev1",
            artifact_id="art1",
            revision_number=1,
            parent_revision_id=None,
            rollback_from_revision_id=None,
            content_ref="item:art1/f.bin",
            inline_content=None,
            size=len(data),
            checksum=checksum,
            mime=mime,
            kind=kind,
            created_at=now,
        )
    return rev, data


def _opts(
    *,
    content_profile: ContentProfile = ContentProfile.TEXT,
    artifact_name: str = "artifact",
) -> types.MappingProxyType[str, object]:
    return types.MappingProxyType({
        "content_profile": content_profile,
        "artifact_name": artifact_name,
    })


@contextmanager
def _assert_no_network():
    """Context manager that counts socket.socket calls and asserts zero."""
    call_count = [0]
    original = socket.socket

    def _counting_socket(*a, **kw):
        call_count[0] += 1
        return original(*a, **kw)

    with patch("socket.socket", _counting_socket):
        yield
    assert call_count[0] == 0, f"network access detected: {call_count[0]} socket calls"


# ---------------------------------------------------------------------------
# Capabilities matrix
# ---------------------------------------------------------------------------


async def test_capabilities_matrix_full():
    ex = OfficeArtifactExporter()

    # markdown/document -> docx, html, original, pptx
    assert await ex.capabilities(
        ArtifactKind.MARKDOWN, "text/markdown", ContentProfile.MARKDOWN
    ) == ("docx", "html", "original", "pptx")
    assert await ex.capabilities(
        ArtifactKind.DOCUMENT, "text/plain", ContentProfile.MARKDOWN
    ) == ("docx", "html", "original", "pptx")

    # TEXT kind -> docx, original
    assert await ex.capabilities(
        ArtifactKind.TEXT, "text/plain", ContentProfile.TEXT
    ) == ("docx", "original")

    # csv + tabular -> original, xlsx
    assert await ex.capabilities(
        ArtifactKind.CSV, "text/csv", ContentProfile.TABULAR
    ) == ("original", "xlsx")

    # json + tabular -> original, xlsx
    assert await ex.capabilities(
        ArtifactKind.JSON, "application/json", ContentProfile.TABULAR
    ) == ("original", "xlsx")

    # data + tabular -> original, xlsx
    assert await ex.capabilities(
        ArtifactKind.DATA, "application/json", ContentProfile.TABULAR
    ) == ("original", "xlsx")

    # json + text -> original
    assert await ex.capabilities(
        ArtifactKind.JSON, "application/json", ContentProfile.TEXT
    ) == ("original",)

    # binary kinds -> original
    for kind in (ArtifactKind.IMAGE, ArtifactKind.PDF, ArtifactKind.OTHER):
        assert await ex.capabilities(
            kind, "application/octet-stream", ContentProfile.BINARY
        ) == ("original",)

    # csv + text (non-tabular) -> original
    assert await ex.capabilities(
        ArtifactKind.CSV, "text/csv", ContentProfile.TEXT
    ) == ("original",)

    # data + text -> original
    assert await ex.capabilities(
        ArtifactKind.DATA, "application/octet-stream", ContentProfile.TEXT
    ) == ("original",)

    # code/html -> original
    assert await ex.capabilities(
        ArtifactKind.CODE, "text/x-python", ContentProfile.TEXT
    ) == ("original",)
    assert await ex.capabilities(
        ArtifactKind.HTML, "text/html", ContentProfile.TEXT
    ) == ("original",)

    # format names are lowercase
    for fmts in [
        await ex.capabilities(ArtifactKind.MARKDOWN, "text/markdown", ContentProfile.MARKDOWN),
        await ex.capabilities(ArtifactKind.CSV, "text/csv", ContentProfile.TABULAR),
    ]:
        assert all(f == f.lower() for f in fmts)


async def test_capabilities_are_deduplicated_and_sorted():
    ex = OfficeArtifactExporter()
    fmts = await ex.capabilities(
        ArtifactKind.MARKDOWN, "text/markdown", ContentProfile.MARKDOWN
    )
    assert len(fmts) == len(set(fmts))
    assert list(fmts) == sorted(fmts)


# ---------------------------------------------------------------------------
# DOCX round-trip
# ---------------------------------------------------------------------------


async def test_docx_roundtrip():
    md = b"# T1\n\nHello world.\n"
    rev, data = _rev(md, kind=ArtifactKind.MARKDOWN, mime="text/markdown")
    ex = OfficeArtifactExporter()
    result = await ex.export(rev, data, "docx", _opts(content_profile=ContentProfile.MARKDOWN))
    assert isinstance(result, ExportedArtifact)
    assert result.mime == _DOCX_MIME
    assert result.filename.endswith(".docx")
    assert result.data

    import docx
    doc = docx.Document(io.BytesIO(result.data))
    assert doc.paragraphs[0].text == "T1"


async def test_docx_format_case_insensitive():
    md = b"# T1\n\nHello.\n"
    rev, data = _rev(md, kind=ArtifactKind.MARKDOWN, mime="text/markdown")
    ex = OfficeArtifactExporter()
    result = await ex.export(rev, data, "DOCX", _opts(content_profile=ContentProfile.MARKDOWN))
    assert result.mime == _DOCX_MIME


# ---------------------------------------------------------------------------
# XLSX formula injection defence
# ---------------------------------------------------------------------------


async def test_xlsx_formula_injection_blocked():
    csv_data = b"col\n=CMD()\n+2\n-3\n@SUM\n"
    rev, data = _rev(csv_data, kind=ArtifactKind.CSV, mime="text/csv")
    ex = OfficeArtifactExporter()
    result = await ex.export(rev, data, "xlsx", _opts(content_profile=ContentProfile.TABULAR))
    assert result.mime == _XLSX_MIME

    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(result.data))
    ws = wb["Artifact"]
    # row 0 = header "col"; rows 1-4 = dangerous values
    dangerous = ["=CMD()", "+2", "-3", "@SUM"]
    for i, expected in enumerate(dangerous, start=1):
        cell = ws.cell(row=i + 1, column=1)
        assert cell.data_type == "s", f"row {i+1} data_type={cell.data_type} expected s"
        assert cell.value == expected


async def test_xlsx_bom_whitespace_formula_injection_and_types():
    # BOM + =, leading whitespace + +, -, @, strict decimal, string
    csv_data = (
        "﻿bom_eq\n"
        "=CMD()\n"
        " +2\n"
        "-3\n"
        "@SUM\n"
        "123\n"
        "45.6\n"
        "hello\n"
    ).encode("utf-8")
    rev, data = _rev(csv_data, kind=ArtifactKind.CSV, mime="text/csv")
    ex = OfficeArtifactExporter()
    result = await ex.export(rev, data, "xlsx", _opts(content_profile=ContentProfile.TABULAR))

    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(result.data))
    ws = wb["Artifact"]
    # header is "bom_eq" (BOM stripped)
    assert ws.cell(row=1, column=1).value == "bom_eq"
    # dangerous values -> string, literal preserved
    assert ws.cell(row=2, column=1).data_type == "s"
    assert ws.cell(row=2, column=1).value == "=CMD()"
    assert ws.cell(row=3, column=1).data_type == "s"
    assert ws.cell(row=3, column=1).value == " +2"
    assert ws.cell(row=4, column=1).data_type == "s"
    assert ws.cell(row=4, column=1).value == "-3"
    assert ws.cell(row=5, column=1).data_type == "s"
    assert ws.cell(row=5, column=1).value == "@SUM"
    # strict decimal -> number
    assert ws.cell(row=6, column=1).data_type == "n"
    assert ws.cell(row=6, column=1).value == 123
    assert ws.cell(row=7, column=1).data_type == "n"
    assert ws.cell(row=7, column=1).value == 45.6
    # plain string -> string
    assert ws.cell(row=8, column=1).data_type == "s"
    assert ws.cell(row=8, column=1).value == "hello"


async def test_xlsx_json_number_and_null_types():
    json_data = b'[{"num": 42, "null_val": null, "neg": -5, "str": "=CMD()"}]'
    rev, data = _rev(json_data, kind=ArtifactKind.JSON, mime="application/json")
    ex = OfficeArtifactExporter()
    result = await ex.export(rev, data, "xlsx", _opts(content_profile=ContentProfile.TABULAR))

    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(result.data))
    ws = wb["Artifact"]
    # header row: num, null_val, neg, str
    assert ws.cell(row=1, column=1).value == "num"
    assert ws.cell(row=1, column=2).value == "null_val"
    assert ws.cell(row=1, column=3).value == "neg"
    assert ws.cell(row=1, column=4).value == "str"
    # data row
    assert ws.cell(row=2, column=1).data_type == "n"
    assert ws.cell(row=2, column=1).value == 42
    # null -> empty cell
    assert ws.cell(row=2, column=2).value is None
    # JSON number -5 -> number (JSON numbers are already typed, no formula check)
    assert ws.cell(row=2, column=3).data_type == "n"
    assert ws.cell(row=2, column=3).value == -5
    # string "=CMD()" -> forced to string (formula injection)
    assert ws.cell(row=2, column=4).data_type == "s"
    assert ws.cell(row=2, column=4).value == "=CMD()"


async def test_xlsx_no_formulas_external_links_macros():
    csv_data = b"a,b\n1,2\n"
    rev, data = _rev(csv_data, kind=ArtifactKind.CSV, mime="text/csv")
    ex = OfficeArtifactExporter()
    result = await ex.export(rev, data, "xlsx", _opts(content_profile=ContentProfile.TABULAR))

    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(result.data))
    ws = wb["Artifact"]
    for row in ws.iter_rows():
        for cell in row:
            if cell.value is not None:
                assert not str(cell.value).startswith("="), "formula found"
    assert wb.vba_archive is None


# ---------------------------------------------------------------------------
# XLSX size limits
# ---------------------------------------------------------------------------


async def test_xlsx_too_large_raises():
    # 100001 rows -> exceeds max_rows
    lines = ["h"] + [f"r{i}" for i in range(100001)]
    csv_data = "\n".join(lines).encode("utf-8")
    rev, data = _rev(csv_data, kind=ArtifactKind.CSV, mime="text/csv")
    ex = OfficeArtifactExporter()
    with pytest.raises(ArtifactExportTooLargeError):
        await ex.export(rev, data, "xlsx", _opts(content_profile=ContentProfile.TABULAR))


async def test_xlsx_too_many_columns_raises():
    header = ",".join(f"c{i}" for i in range(257))
    csv_data = (header + "\n").encode("utf-8")
    rev, data = _rev(csv_data, kind=ArtifactKind.CSV, mime="text/csv")
    ex = OfficeArtifactExporter()
    with pytest.raises(ArtifactExportTooLargeError):
        await ex.export(rev, data, "xlsx", _opts(content_profile=ContentProfile.TABULAR))


# ---------------------------------------------------------------------------
# Office output metadata scrubbing
# ---------------------------------------------------------------------------


def _assert_no_local_path_in_zip(data: bytes) -> None:
    with zipfile.ZipFile(io.BytesIO(data)) as zf:
        for name in zf.namelist():
            content = zf.read(name)
            assert b"/Users/" not in content, f"local path found in {name}"


async def test_office_output_has_no_local_metadata():
    md = b"# Title\n\nSome text.\n"
    rev, data = _rev(md, kind=ArtifactKind.MARKDOWN, mime="text/markdown")
    ex = OfficeArtifactExporter()
    opts = _opts(content_profile=ContentProfile.MARKDOWN, artifact_name="test")

    # DOCX
    docx_result = await ex.export(rev, data, "docx", opts)
    import docx
    doc = docx.Document(io.BytesIO(docx_result.data))
    cp = doc.core_properties
    assert not cp.author
    assert not cp.last_modified_by
    _assert_no_local_path_in_zip(docx_result.data)

    # PPTX
    pptx_result = await ex.export(rev, data, "pptx", opts)
    import pptx
    prs = pptx.Presentation(io.BytesIO(pptx_result.data))
    pcp = prs.core_properties
    assert not pcp.author
    assert not pcp.last_modified_by
    _assert_no_local_path_in_zip(pptx_result.data)

    # XLSX
    csv_data = b"a,b\n1,2\n"
    csv_rev, csv_bytes = _rev(csv_data, kind=ArtifactKind.CSV, mime="text/csv")
    xlsx_result = await ex.export(
        csv_rev, csv_bytes, "xlsx", _opts(content_profile=ContentProfile.TABULAR)
    )
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(xlsx_result.data))
    assert not wb.properties.creator
    assert not wb.properties.lastModifiedBy
    _assert_no_local_path_in_zip(xlsx_result.data)


# ---------------------------------------------------------------------------
# Injectable limits: never truncate, always raise
# ---------------------------------------------------------------------------


async def test_office_limits_are_injectable_and_never_truncate():
    # max_blocks: markdown with too many paragraphs
    md = "\n\n".join(f"paragraph {i}" for i in range(10)).encode()
    rev, data = _rev(md, kind=ArtifactKind.MARKDOWN, mime="text/markdown")
    ex = OfficeArtifactExporter(config=ArtifactExporterConfig(max_blocks=5))
    with pytest.raises(ArtifactExportTooLargeError):
        await ex.export(rev, data, "docx", _opts(content_profile=ContentProfile.MARKDOWN))

    # max_table_cells: table with too many cells
    md_table = "| a | b | c |\n|---|---|---|\n| 1 | 2 | 3 |\n| 4 | 5 | 6 |\n"
    rev2, data2 = _rev(md_table.encode(), kind=ArtifactKind.MARKDOWN, mime="text/markdown")
    ex2 = OfficeArtifactExporter(config=ArtifactExporterConfig(max_table_cells=4))
    with pytest.raises(ArtifactExportTooLargeError):
        await ex2.export(rev2, data2, "docx", _opts(content_profile=ContentProfile.MARKDOWN))

    # max_slides: too many H2 sections for pptx
    md_slides = "# Title\n\n" + "\n\n".join(f"## Slide {i}\n\ntext" for i in range(5))
    rev3, data3 = _rev(md_slides.encode(), kind=ArtifactKind.MARKDOWN, mime="text/markdown")
    ex3 = OfficeArtifactExporter(config=ArtifactExporterConfig(max_slides=2))
    with pytest.raises(ArtifactExportTooLargeError):
        await ex3.export(rev3, data3, "pptx", _opts(content_profile=ContentProfile.MARKDOWN))

    # max_rows: CSV with too many rows
    csv_rows = "h\n" + "\n".join(f"r{i}" for i in range(10))
    rev4, data4 = _rev(csv_rows.encode(), kind=ArtifactKind.CSV, mime="text/csv")
    ex4 = OfficeArtifactExporter(config=ArtifactExporterConfig(max_rows=3))
    with pytest.raises(ArtifactExportTooLargeError):
        await ex4.export(rev4, data4, "xlsx", _opts(content_profile=ContentProfile.TABULAR))

    # max_columns: CSV with too many columns
    header = ",".join(f"c{i}" for i in range(5))
    csv_cols = header + "\n"
    rev5, data5 = _rev(csv_cols.encode(), kind=ArtifactKind.CSV, mime="text/csv")
    ex5 = OfficeArtifactExporter(config=ArtifactExporterConfig(max_columns=3))
    with pytest.raises(ArtifactExportTooLargeError):
        await ex5.export(rev5, data5, "xlsx", _opts(content_profile=ContentProfile.TABULAR))

    # max_cell_chars: cell too long
    long_cell = "x" * 100
    csv_long = f"h\n{long_cell}\n"
    rev6, data6 = _rev(csv_long.encode(), kind=ArtifactKind.CSV, mime="text/csv")
    ex6 = OfficeArtifactExporter(config=ArtifactExporterConfig(max_cell_chars=10))
    with pytest.raises(ArtifactExportTooLargeError):
        await ex6.export(rev6, data6, "xlsx", _opts(content_profile=ContentProfile.TABULAR))

    # max_output_bytes: output too large
    big_md = ("# H\n\n" + "paragraph\n\n" * 100).encode()
    rev7, data7 = _rev(big_md, kind=ArtifactKind.MARKDOWN, mime="text/markdown")
    ex7 = OfficeArtifactExporter(config=ArtifactExporterConfig(max_output_bytes=100))
    with pytest.raises(ArtifactExportTooLargeError):
        await ex7.export(rev7, data7, "docx", _opts(content_profile=ContentProfile.MARKDOWN))


# ---------------------------------------------------------------------------
# Markdown structure: DOCX + PPTX + unsafe nodes
# ---------------------------------------------------------------------------


_MD_STRUCTURE = (
    "# Main Title\n\n"
    "## Section One\n\n"
    "A paragraph here.\n\n"
    "- item one\n- item two\n\n"
    "1. first\n2. second\n\n"
    "```python\nprint('hello')\n```\n\n"
    "| Col A | Col B |\n|-------|-------|\n| 1 | 2 |\n| 3 | 4 |\n"
)


async def test_docx_markdown_structure():
    with _assert_no_network():
        rev, data = _rev(_MD_STRUCTURE.encode(), kind=ArtifactKind.MARKDOWN, mime="text/markdown")
        ex = OfficeArtifactExporter()
        result = await ex.export(rev, data, "docx", _opts(content_profile=ContentProfile.MARKDOWN))

    import docx
    doc = docx.Document(io.BytesIO(result.data))
    texts = [p.text for p in doc.paragraphs if p.text.strip()]
    # headings preserved
    assert "Main Title" in texts
    assert "Section One" in texts
    # paragraph preserved
    assert "A paragraph here." in texts
    # list items preserved
    assert "item one" in texts
    assert "item two" in texts
    assert "first" in texts
    assert "second" in texts
    # code block content preserved
    assert "print('hello')" in texts
    # table preserved
    assert len(doc.tables) >= 1
    tbl = doc.tables[0]
    assert tbl.rows[0].cells[0].text == "Col A"
    assert tbl.rows[0].cells[1].text == "Col B"
    assert tbl.rows[1].cells[0].text == "1"
    assert tbl.rows[1].cells[1].text == "2"


async def test_docx_headings_all_levels():
    md = "# H1\n\n## H2\n\n### H3\n\n#### H4\n\n##### H5\n\n###### H6\n"
    rev, data = _rev(md.encode(), kind=ArtifactKind.MARKDOWN, mime="text/markdown")
    ex = OfficeArtifactExporter()
    result = await ex.export(rev, data, "docx", _opts(content_profile=ContentProfile.MARKDOWN))

    import docx
    doc = docx.Document(io.BytesIO(result.data))
    headings = [p.text for p in doc.paragraphs if p.text.strip()]
    for level in range(1, 7):
        assert f"H{level}" in headings, f"H{level} not found in {headings}"


async def test_pptx_markdown_structure():
    with _assert_no_network():
        rev, data = _rev(_MD_STRUCTURE.encode(), kind=ArtifactKind.MARKDOWN, mime="text/markdown")
        ex = OfficeArtifactExporter()
        result = await ex.export(rev, data, "pptx", _opts(content_profile=ContentProfile.MARKDOWN))

    import pptx
    prs = pptx.Presentation(io.BytesIO(result.data))
    assert len(prs.slides) >= 1
    # collect all text from slides
    all_text: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text.append(shape.text_frame.text)
    joined = "\n".join(all_text)
    assert "Main Title" in joined
    assert "Section One" in joined
    assert "A paragraph here." in joined


async def test_docx_inline_html_degrades_to_text():
    md = "# Title\n\n<script>alert(1)</script>\n\n<div>not rendered</div>\n"
    rev, data = _rev(md.encode(), kind=ArtifactKind.MARKDOWN, mime="text/markdown")
    ex = OfficeArtifactExporter()
    result = await ex.export(rev, data, "docx", _opts(content_profile=ContentProfile.MARKDOWN))

    import docx
    doc = docx.Document(io.BytesIO(result.data))
    full_text = " ".join(p.text for p in doc.paragraphs)
    # HTML tags must not appear as live tags; content degrades to text
    assert "<script>" not in full_text
    assert "<div>" not in full_text
    # the text content is preserved
    assert "alert(1)" in full_text


async def test_docx_image_reference_is_placeholder():
    with _assert_no_network():
        md = "# Title\n\n![alt text](http://example.com/img.png)\n"
        rev, data = _rev(md.encode(), kind=ArtifactKind.MARKDOWN, mime="text/markdown")
        ex = OfficeArtifactExporter()
        result = await ex.export(rev, data, "docx", _opts(content_profile=ContentProfile.MARKDOWN))

    import docx
    doc = docx.Document(io.BytesIO(result.data))
    full_text = " ".join(p.text for p in doc.paragraphs)
    # no image embedded (v1 degrades to placeholder text)
    assert "alt text" in full_text
    # check no image relationships in the docx
    with zipfile.ZipFile(io.BytesIO(result.data)) as zf:
        has_image = any(
            "image" in name or "media" in name
            for name in zf.namelist()
        )
    assert not has_image


async def test_pptx_image_reference_is_placeholder():
    with _assert_no_network():
        md = "# Title\n\n![pic](http://example.com/img.png)\n"
        rev, data = _rev(md.encode(), kind=ArtifactKind.MARKDOWN, mime="text/markdown")
        ex = OfficeArtifactExporter()
        result = await ex.export(rev, data, "pptx", _opts(content_profile=ContentProfile.MARKDOWN))

    with zipfile.ZipFile(io.BytesIO(result.data)) as zf:
        has_image = any(
            "image" in name or "media" in name
            for name in zf.namelist()
        )
    assert not has_image


# ---------------------------------------------------------------------------
# TEXT kind -> DOCX (plain paragraphs)
# ---------------------------------------------------------------------------


async def test_text_kind_docx_export():
    rev, data = _rev(b"plain text line\nsecond line", kind=ArtifactKind.TEXT, mime="text/plain")
    ex = OfficeArtifactExporter()
    result = await ex.export(rev, data, "docx", _opts(content_profile=ContentProfile.TEXT))
    assert result.mime == _DOCX_MIME

    import docx
    doc = docx.Document(io.BytesIO(result.data))
    texts = [p.text for p in doc.paragraphs if p.text.strip()]
    assert "plain text line" in texts
    assert "second line" in texts


# ---------------------------------------------------------------------------
# HTML export
# ---------------------------------------------------------------------------


async def test_html_export_from_markdown():
    md = b"# Title\n\n[link](http://e.com)\n"
    rev, data = _rev(md, kind=ArtifactKind.MARKDOWN, mime="text/markdown")
    ex = OfficeArtifactExporter()
    result = await ex.export(rev, data, "html", _opts(content_profile=ContentProfile.MARKDOWN))
    assert result.mime == _HTML_MIME
    assert result.filename.endswith(".html")
    html = result.data.decode("utf-8")
    assert "<h1>Title</h1>" in html
    assert 'href="http://e.com"' in html
    assert "Content-Security-Policy" in html


# ---------------------------------------------------------------------------
# Original export
# ---------------------------------------------------------------------------


async def test_original_export_returns_raw_bytes():
    raw = b"\x89PNG\r\n\x1a\n\x00\x00"
    rev, data = _rev(raw, kind=ArtifactKind.IMAGE, mime="image/png")
    ex = OfficeArtifactExporter()
    result = await ex.export(rev, data, "original", _opts(content_profile=ContentProfile.BINARY, artifact_name="photo.png"))
    assert result.data == raw
    assert result.mime == "image/png"
    assert result.filename == "photo.png"


async def test_original_export_text_kind():
    rev, data = _rev(b"hello world", kind=ArtifactKind.TEXT, mime="text/plain")
    ex = OfficeArtifactExporter()
    result = await ex.export(rev, data, "original", _opts(content_profile=ContentProfile.TEXT, artifact_name="notes.txt"))
    assert result.data == b"hello world"
    assert result.mime == "text/plain"


# ---------------------------------------------------------------------------
# Unsupported format
# ---------------------------------------------------------------------------


async def test_unsupported_format_raises():
    rev, data = _rev(b"# test", kind=ArtifactKind.MARKDOWN, mime="text/markdown")
    ex = OfficeArtifactExporter()
    with pytest.raises(ArtifactExportUnsupportedError):
        await ex.export(rev, data, "pdf", _opts(content_profile=ContentProfile.MARKDOWN))


async def test_binary_kind_docx_not_supported():
    rev, data = _rev(b"\x00\x01", kind=ArtifactKind.IMAGE, mime="image/png")
    ex = OfficeArtifactExporter()
    with pytest.raises(ArtifactExportUnsupportedError):
        await ex.export(rev, data, "docx", _opts(content_profile=ContentProfile.BINARY))


# ---------------------------------------------------------------------------
# XLSX from JSON array-of-arrays
# ---------------------------------------------------------------------------


async def test_xlsx_from_json_array_of_arrays():
    json_data = b'[[1, 2, 3], [4, 5, 6]]'
    rev, data = _rev(json_data, kind=ArtifactKind.JSON, mime="application/json")
    ex = OfficeArtifactExporter()
    result = await ex.export(rev, data, "xlsx", _opts(content_profile=ContentProfile.TABULAR))

    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(result.data))
    ws = wb["Artifact"]
    assert ws.cell(row=1, column=1).value == 1
    assert ws.cell(row=1, column=2).value == 2
    assert ws.cell(row=2, column=3).value == 6


async def test_xlsx_from_json_array_of_objects():
    json_data = b'[{"name": "Alice", "age": 30}, {"name": "Bob", "age": 25}]'
    rev, data = _rev(json_data, kind=ArtifactKind.JSON, mime="application/json")
    ex = OfficeArtifactExporter()
    result = await ex.export(rev, data, "xlsx", _opts(content_profile=ContentProfile.TABULAR))

    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(result.data))
    ws = wb["Artifact"]
    # header row uses object keys
    assert ws.cell(row=1, column=1).value == "name"
    assert ws.cell(row=1, column=2).value == "age"
    # data rows
    assert ws.cell(row=2, column=1).value == "Alice"
    assert ws.cell(row=2, column=2).value == 30
    assert ws.cell(row=3, column=1).value == "Bob"
    assert ws.cell(row=3, column=2).value == 25
